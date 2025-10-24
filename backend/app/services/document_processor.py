"""
Document processing service for EduScribe backend.
Based on your existing ingest_corpus.py
"""
import os
import faiss
import pickle
import asyncio
from pathlib import Path
from typing import List, Dict, Any
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
from pptx import Presentation
import docx

from app.core.config import settings

# Global embedder (lazy loaded)
_embedder = None

def get_embedder():
    """Get or create the sentence transformer model."""
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer(settings.EMBEDDING_MODEL)
    return _embedder

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from PDF file."""
    try:
        reader = PdfReader(pdf_path)
        text = []
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text.append(content)
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting PDF {pdf_path}: {e}")
        return ""

def extract_text_from_ppt(ppt_path: str) -> str:
    """Extract text from PPTX file."""
    try:
        prs = Presentation(ppt_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text.append(shape.text.strip())
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting PPT {ppt_path}: {e}")
        return ""

def extract_text_from_docx(docx_path: str) -> str:
    """Extract text from DOCX file."""
    try:
        doc = docx.Document(docx_path)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text.strip())
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting DOCX {docx_path}: {e}")
        return ""

def extract_text_from_txt(txt_path: str) -> str:
    """Extract text from TXT file."""
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"Error extracting TXT {txt_path}: {e}")
        return ""

def chunk_text(text: str, chunk_size: int = 300) -> List[str]:
    """Split text into word chunks."""
    words = text.split()
    return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def extract_text_from_file(file_path: str) -> str:
    """Extract text from various file formats."""
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()
    
    if suffix == ".pdf":
        return extract_text_from_pdf(str(file_path))
    elif suffix == ".pptx":
        return extract_text_from_ppt(str(file_path))
    elif suffix in [".docx", ".doc"]:
        return extract_text_from_docx(str(file_path))
    elif suffix == ".txt":
        return extract_text_from_txt(str(file_path))
    else:
        print(f"Unsupported file format: {suffix}")
        return ""

async def build_faiss_for_lecture(lecture_id: str, document_paths: List[str]) -> Dict[str, Any]:
    """
    Build FAISS index for a specific lecture from its documents.
    """
    lecture_dir = Path(settings.PROCESSED_DIR) / lecture_id
    lecture_dir.mkdir(parents=True, exist_ok=True)
    
    index_file = lecture_dir / "vector_index.faiss"
    chunks_file = lecture_dir / "chunks.pkl"
    metadata_file = lecture_dir / "metadata.json"

    # Run in executor to avoid blocking
    loop = asyncio.get_event_loop()
    result = await loop.run_in_executor(
        None, 
        _build_faiss_sync, 
        lecture_id, 
        document_paths, 
        str(index_file), 
        str(chunks_file),
        str(metadata_file)
    )
    
    return result

def _build_faiss_sync(lecture_id: str, document_paths: List[str], index_file: str, chunks_file: str, metadata_file: str) -> Dict[str, Any]:
    """Synchronous FAISS building function."""
    embedder = get_embedder()
    
    all_chunks = []
    chunk_metadata = []
    
    for doc_path in document_paths:
        try:
            print(f"Processing document: {doc_path}")
            text = extract_text_from_file(doc_path)
            
            if not text.strip():
                continue
                
            # Chunk the text
            doc_chunks = chunk_text(text, chunk_size=300)
            
            # Add metadata for each chunk
            for i, chunk in enumerate(doc_chunks):
                all_chunks.append(chunk)
                chunk_metadata.append({
                    "document_path": doc_path,
                    "chunk_index": i,
                    "document_name": Path(doc_path).name
                })
                
        except Exception as e:
            print(f"Error processing document {doc_path}: {e}")
            continue

    if not all_chunks:
        return {"error": "No chunks extracted", "chunks_count": 0}

    print(f"Total chunks for lecture {lecture_id}: {len(all_chunks)}")

    # Create embeddings
    embeddings = embedder.encode(all_chunks, convert_to_numpy=True)

    # Build FAISS index
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings)

    # Save files
    faiss.write_index(index, index_file)
    
    with open(chunks_file, "wb") as f:
        pickle.dump(all_chunks, f)
    
    # Save metadata as JSON
    import json
    metadata = {
        "lecture_id": lecture_id,
        "chunks_count": len(all_chunks),
        "documents": document_paths,
        "chunk_metadata": chunk_metadata
    }
    
    with open(metadata_file, "w") as f:
        json.dump(metadata, f, indent=2)

    print(f"✅ FAISS index built for lecture {lecture_id}")
    
    return {
        "success": True,
        "chunks_count": len(all_chunks),
        "documents_processed": len(document_paths),
        "index_file": index_file,
        "chunks_file": chunks_file
    }

def load_faiss_index(lecture_id: str) -> tuple:
    """Load FAISS index and chunks for a lecture."""
    lecture_dir = Path(settings.PROCESSED_DIR) / lecture_id
    index_file = lecture_dir / "vector_index.faiss"
    chunks_file = lecture_dir / "chunks.pkl"
    
    if not index_file.exists() or not chunks_file.exists():
        return None, None
    
    try:
        index = faiss.read_index(str(index_file))
        with open(chunks_file, "rb") as f:
            chunks = pickle.load(f)
        return index, chunks
    except Exception as e:
        print(f"Error loading FAISS index for lecture {lecture_id}: {e}")
        return None, None

def query_documents(query_text: str, lecture_id: str, top_k: int = 3) -> List[str]:
    """Query the FAISS index for relevant document chunks."""
    index, chunks = load_faiss_index(lecture_id)
    
    if index is None or chunks is None:
        return []
    
    try:
        embedder = get_embedder()
        query_embedding = embedder.encode([query_text], convert_to_numpy=True)
        
        distances, indices = index.search(query_embedding, top_k)
        
        results = []
        for i in indices[0]:
            if i < len(chunks):
                results.append(chunks[i])
        
        return results
    except Exception as e:
        print(f"Error querying documents for lecture {lecture_id}: {e}")
        return []
